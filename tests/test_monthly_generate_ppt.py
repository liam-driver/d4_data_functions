import sys
import os
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

import json
import shutil

import pytest
from pptx import Presentation
from pptx.util import Inches

import monthly_reports.generate_ppt as generate_ppt_module
from monthly_reports.generate_ppt import (
    C,
    STATUS_COLOURS,
    SLD_LAYOUT_BLANK,
    SLD_LAYOUT_SECTION_SEPARATOR,
    SLD_LAYOUT_SUNSET_SECTION_SEPARATOR,
    SLD_LAYOUT_TITLE_AND_BODY,
    _extract_current_tasks,
    _add_table_shape,
    _fmt_date,
    _assemble_pptx,
    generate_skeleton_ppt,
    generate_ppt,
    slide_planning_gantt,
    slide_scorecard_commentary,
)


@pytest.fixture
def prs():
    p = Presentation('slides/template.pptx')
    slide_rid = p.slides._sldIdLst[0].get(
        '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
    )
    p.part.drop_rel(slide_rid)
    p.slides._sldIdLst.remove(p.slides._sldIdLst[0])
    return p


# ── _extract_current_tasks ────────────────────────────────────────────────────

class TestExtractCurrentTasks:
    def test_returns_tasks_for_current_quarter(self):
        task = {"name": "Current task", "category": "Active Workstream",
                "start_date": "01/01/26", "end_date": "31/12/26"}
        plan = {
            "Q1 2026": {"plan_status": "old",     "tasks": [{"name": "Old task"}]},
            "Q2 2026": {"plan_status": "current",  "tasks": [task]},
        }
        assert _extract_current_tasks(plan) == [task]

    def test_returns_empty_when_no_current_quarter(self):
        plan = {"Q1 2026": {"plan_status": "old", "tasks": [{"name": "Old task"}]}}
        assert _extract_current_tasks(plan) == []

    def test_returns_empty_for_none(self):
        assert _extract_current_tasks(None) == []

    def test_returns_empty_for_empty_dict(self):
        assert _extract_current_tasks({}) == []

    def test_returns_empty_tasks_list_when_current_quarter_has_no_tasks(self):
        plan = {"Q2 2026": {"plan_status": "current", "tasks": []}}
        assert _extract_current_tasks(plan) == []

    def test_handles_flat_list_shape(self):
        tasks = [{"name": "Task A"}, {"name": "Task B"}]
        assert _extract_current_tasks(tasks) == tasks

    def test_handles_tasks_keyed_shape(self):
        plan = {"tasks": [{"name": "Task A"}]}
        assert _extract_current_tasks(plan) == [{"name": "Task A"}]

    def test_picks_current_among_multiple_quarters(self):
        task = {"name": "Current", "category": "Active Workstream",
                "start_date": "01/01/26", "end_date": "31/12/26"}
        plan = {
            "Q4 2025": {"plan_status": "old",     "tasks": [{"name": "Old"}]},
            "Q1 2026": {"plan_status": "old",     "tasks": [{"name": "Also old"}]},
            "Q2 2026": {"plan_status": "current",  "tasks": [task]},
        }
        result = _extract_current_tasks(plan)
        assert len(result) == 1
        assert result[0]["name"] == "Current"

    def test_filters_out_bau_and_undated_tasks(self):
        """Only Active Workstream tasks overlapping the reporting month are returned."""
        keep = {"name": "Keep", "category": "Active Workstream",
                "start_date": "01/01/26", "end_date": "31/12/26"}
        bau = {"name": "BAU task", "category": "BAU",
               "start_date": "01/01/26", "end_date": "31/12/26"}
        undated = {"name": "No dates", "category": "Active Workstream"}
        plan = {"Q2 2026": {"plan_status": "current", "tasks": [keep, bau, undated]}}
        assert _extract_current_tasks(plan) == [keep]


# ── _fmt_date ─────────────────────────────────────────────────────────────────

class TestFmtDate:
    def test_two_digit_year(self):
        assert _fmt_date("06/04/26") == "06/04/2026"

    def test_four_digit_year(self):
        assert _fmt_date("06/04/2026") == "06/04/2026"

    def test_iso_format(self):
        assert _fmt_date("2026-04-06") == "06/04/2026"

    def test_empty_string(self):
        assert _fmt_date("") == ""

    def test_none_returns_empty(self):
        assert _fmt_date(None) == ""


# ── _add_table_shape ──────────────────────────────────────────────────────────

class TestAddTableShape:
    def test_creates_table_with_correct_row_count(self, prs):
        slide = prs.slides.add_slide(prs.slide_layouts[SLD_LAYOUT_BLANK])
        headers = ['Name', 'Status']
        rows    = [['Task A', 'Complete'], ['Task B', 'Scheduled']]
        _add_table_shape(slide, headers, rows, Inches(0.5), Inches(1), Inches(9), Inches(4))

        tables = [s for s in slide.shapes if s.has_table]
        assert len(tables) == 1
        tbl = tables[0].table
        assert len(tbl.rows) == 3  # 1 header + 2 data rows

    def test_header_cell_text(self, prs):
        slide = prs.slides.add_slide(prs.slide_layouts[SLD_LAYOUT_BLANK])
        _add_table_shape(slide, ['Name', 'Status'], [['Task A', 'Complete']],
                         Inches(0.5), Inches(1), Inches(9), Inches(2))
        tbl = [s for s in slide.shapes if s.has_table][0].table
        assert tbl.rows[0].cells[0].text == 'Name'
        assert tbl.rows[0].cells[1].text == 'Status'

    def test_data_cell_text(self, prs):
        slide = prs.slides.add_slide(prs.slide_layouts[SLD_LAYOUT_BLANK])
        _add_table_shape(slide, ['Name'], [['Task A']], Inches(0.5), Inches(1), Inches(9), Inches(2))
        tbl = [s for s in slide.shapes if s.has_table][0].table
        assert tbl.rows[1].cells[0].text == 'Task A'

    def test_header_uses_dark_background(self, prs):
        slide = prs.slides.add_slide(prs.slide_layouts[SLD_LAYOUT_BLANK])
        _add_table_shape(slide, ['Name'], [['Task A']], Inches(0.5), Inches(1), Inches(9), Inches(2))
        tbl = [s for s in slide.shapes if s.has_table][0].table
        assert tbl.cell(0, 0).fill.fore_color.rgb == C["dark"]

    def test_status_col_applies_status_colour_for_complete(self, prs):
        slide = prs.slides.add_slide(prs.slide_layouts[SLD_LAYOUT_BLANK])
        _add_table_shape(slide, ['Name', 'Status'], [['Task A', 'Complete']],
                         Inches(0.5), Inches(1), Inches(9), Inches(2), status_col=1)
        tbl = [s for s in slide.shapes if s.has_table][0].table
        assert tbl.cell(1, 1).fill.fore_color.rgb == STATUS_COLOURS['Complete']

    def test_status_col_applies_status_colour_for_blocked(self, prs):
        slide = prs.slides.add_slide(prs.slide_layouts[SLD_LAYOUT_BLANK])
        _add_table_shape(slide, ['Name', 'Status'], [['Task B', 'Blocked']],
                         Inches(0.5), Inches(1), Inches(9), Inches(2), status_col=1)
        tbl = [s for s in slide.shapes if s.has_table][0].table
        assert tbl.cell(1, 1).fill.fore_color.rgb == STATUS_COLOURS['Blocked']

    def test_alternating_row_shading(self, prs):
        slide = prs.slides.add_slide(prs.slide_layouts[SLD_LAYOUT_BLANK])
        rows = [['A', 'Complete'], ['B', 'Complete'], ['C', 'Complete']]
        _add_table_shape(slide, ['Name', 'Status'], rows,
                         Inches(0.5), Inches(1), Inches(9), Inches(4))
        tbl = [s for s in slide.shapes if s.has_table][0].table
        # Row 1 (i=0) → light; row 2 (i=1) → white; row 3 (i=2) → light
        assert tbl.cell(1, 0).fill.fore_color.rgb == C["light"]
        assert tbl.cell(2, 0).fill.fore_color.rgb == C["white"]
        assert tbl.cell(3, 0).fill.fore_color.rgb == C["light"]


# ── slide_planning_gantt ──────────────────────────────────────────────────────

class TestSlidePlanningGantt:
    TASKS = [
        {'name': 'Task A', 'platform': 'Google Ads',
         'start_date': '06/04/26', 'end_date': '10/04/26', 'status': 'Complete'},
        {'name': 'Task B', 'platform': 'Meta Ads',
         'start_date': '12/04/26', 'end_date': '25/04/26', 'status': 'Scheduled'},
    ]

    def test_returns_none_for_empty_tasks(self, prs):
        assert slide_planning_gantt(prs, 'Q2 Plan', []) is None

    def test_returns_slide_for_valid_tasks(self, prs):
        slide = slide_planning_gantt(prs, 'Q2 Plan', self.TASKS)
        assert slide is not None

    def test_slide_uses_shapes_not_picture(self, prs):
        slide = slide_planning_gantt(prs, 'Q2 Plan', self.TASKS)
        pictures    = [s for s in slide.shapes if s.shape_type == 13]  # PICTURE
        tables      = [s for s in slide.shapes if s.has_table]
        auto_shapes = [s for s in slide.shapes if s.shape_type == 1]   # AUTO_SHAPE
        assert len(pictures) == 0
        assert len(tables) == 0
        assert len(auto_shapes) >= 1

    def test_labels_combine_platform_and_name(self, prs):
        slide = slide_planning_gantt(prs, 'Q2 Plan', self.TASKS)
        all_text = ' '.join(
            run.text
            for shp in slide.shapes if shp.has_text_frame
            for p in shp.text_frame.paragraphs
            for run in p.runs
        )
        assert 'Google Ads: Task A' in all_text
        assert 'Meta Ads: Task B' in all_text


# ── slide_scorecard_commentary ────────────────────────────────────────────────

class TestSlideScorecardCommentary:
    def _make_kpis(self):
        return [
            ('Cost',    {'curr': '£1,000', 'prev': '£900',  'pct': '+11%'}),
            ('Revenue', {'curr': '£5,000', 'prev': '£4,000', 'pct': '+25%'}),
            ('ROAS',    {'curr': '5.0',    'prev': '4.5',    'pct': '+11%'}),
        ]

    def test_title_placeholder_is_set(self, prs):
        slide = slide_scorecard_commentary(
            prs, 'Top Level View', 'Good month', [{'point': 'Revenue up'}], self._make_kpis()
        )
        assert slide.placeholders[0].text == 'Top Level View'

    def test_kpi_boxes_are_added(self, prs):
        kpis = self._make_kpis()
        slide = slide_scorecard_commentary(
            prs, 'Top Level View', 'Good month', [{'point': 'Revenue up'}], kpis
        )
        # Shapes include layout placeholders + 3 KPI boxes; at least 3 shapes total
        assert len(slide.shapes) >= len(kpis)

    def test_returns_slide_object(self, prs):
        slide = slide_scorecard_commentary(
            prs, 'Top Level View', 'Good month', [{'point': 'Revenue up'}], self._make_kpis()
        )
        assert slide is not None


# ── _assemble_pptx (Team-based skeleton/final assembly, ADR 0012) ─────────────

def _plan_json(task_name):
    return {
        "Q1 2026": {
            "plan_status": "current",
            "plan_start": "01/01/26",
            "plan_end": "31/03/26",
            "tasks": [
                {"name": task_name, "category": "Active Workstream", "platform": "Google Ads",
                 "start_date": "01/01/26", "end_date": "31/01/26", "status": "Scheduled"},
            ],
        }
    }


def _separator_titles(prs, layout_idx):
    layout = prs.slide_layouts[layout_idx]
    return [slide.placeholders[0].text for slide in prs.slides if slide.slide_layout == layout]


def _all_slide_text(prs):
    return ' '.join(
        run.text
        for slide in prs.slides
        for shp in slide.shapes if shp.has_text_frame
        for p in shp.text_frame.paragraphs
        for run in p.runs
    )


class TestAssemblePptxTeams:
    def _client(self):
        return {
            "name": "Acme",
            "start_date_string": "01/04/2026",
            "end_date_string": "30/04/2026",
            "account_type": "Ecommerce",
            "paid_data_mom": {"Total": {
                "Cost":               {"curr": "£1,000", "prev": "£900",   "pct": "+11%"},
                "Transaction Revenue": {"curr": "£5,000", "prev": "£4,000", "pct": "+25%"},
                "ROAS":               {"curr": "5.0",    "prev": "4.5",   "pct": "+11%"},
                "Conversion Rate":    {"curr": "2%",     "prev": "1.8%",  "pct": "+11%"},
            }},
        }

    def _teams_content(self, include_trends=False):
        ppc_block = {
            "team": "ppc",
            "overviews": [{"data_key": "paid_data", "comparison": "mom"}],
            "plan_json": _plan_json("PPC Task"),
        }
        if include_trends:
            ppc_block["trends"] = [{
                "title": "Paid Search Grows", "summary": "ROAS improves",
                "bullets": [{"point": "ROAS up"}], "graph": {},
            }]
        seo_block = {"team": "seo", "overviews": [], "plan_json": _plan_json("SEO Task")}
        cro_block = {"team": "cro", "overviews": [], "plan_json": _plan_json("CRO Task")}
        return {"teams": [ppc_block, seo_block, cro_block]}

    def test_renders_navy_section_per_active_team_in_order(self, tmp_path):
        output_path = str(tmp_path / "draft.pptx")
        _assemble_pptx(self._client(), self._teams_content(), output_path)
        prs = Presentation(output_path)
        assert _separator_titles(prs, SLD_LAYOUT_SECTION_SEPARATOR) == ['Paid Media', 'SEO', 'CRO']

    def test_skips_team_with_no_block(self, tmp_path):
        teams_content = self._teams_content()
        teams_content['teams'] = [t for t in teams_content['teams'] if t['team'] != 'seo']
        output_path = str(tmp_path / "draft.pptx")
        _assemble_pptx(self._client(), teams_content, output_path)
        prs = Presentation(output_path)
        assert _separator_titles(prs, SLD_LAYOUT_SECTION_SEPARATOR) == ['Paid Media', 'CRO']

    def test_top_level_trends_omitted_when_ppc_has_no_trends(self, tmp_path):
        output_path = str(tmp_path / "draft.pptx")
        _assemble_pptx(self._client(), self._teams_content(include_trends=False), output_path)
        prs = Presentation(output_path)
        assert 'Top Level Trends' not in _separator_titles(prs, SLD_LAYOUT_SUNSET_SECTION_SEPARATOR)

    def test_top_level_trends_rendered_when_ppc_has_trends(self, tmp_path):
        output_path = str(tmp_path / "final.pptx")
        _assemble_pptx(self._client(), self._teams_content(include_trends=True), output_path)
        prs = Presentation(output_path)
        assert 'Top Level Trends' in _separator_titles(prs, SLD_LAYOUT_SUNSET_SECTION_SEPARATOR)

    def test_only_ppc_ever_renders_trends_section(self, tmp_path):
        """SEO/CRO never get a Top Level Trends section, even if trends were injected there."""
        teams_content = self._teams_content()
        teams_content['teams'][1]['trends'] = [
            {"title": "x", "summary": "y", "bullets": [], "graph": {}}
        ]
        output_path = str(tmp_path / "draft.pptx")
        _assemble_pptx(self._client(), teams_content, output_path)
        prs = Presentation(output_path)
        assert 'Top Level Trends' not in _separator_titles(prs, SLD_LAYOUT_SUNSET_SECTION_SEPARATOR)

    def test_each_team_gets_its_own_kanban_and_gantt(self, tmp_path):
        output_path = str(tmp_path / "draft.pptx")
        _assemble_pptx(self._client(), self._teams_content(), output_path)
        prs = Presentation(output_path)
        all_text = _all_slide_text(prs)
        assert 'Google Ads: PPC Task' in all_text
        assert 'Google Ads: SEO Task' in all_text
        assert 'Google Ads: CRO Task' in all_text

    def test_team_with_no_plan_tasks_gets_no_plan_overview_section(self, tmp_path):
        teams_content = {"teams": [{"team": "ppc", "overviews": [], "plan_json": None}]}
        output_path = str(tmp_path / "draft.pptx")
        _assemble_pptx(self._client(), teams_content, output_path)
        prs = Presentation(output_path)
        assert 'Plan Overview' not in _separator_titles(prs, SLD_LAYOUT_SUNSET_SECTION_SEPARATOR)


# ── generate_skeleton_ppt / generate_ppt orchestration (ADR 0012) ─────────────

@pytest.fixture
def isolated_project(tmp_path, monkeypatch):
    (tmp_path / "storage").mkdir()
    (tmp_path / "slides").mkdir()
    shutil.copy("slides/template.pptx", tmp_path / "slides" / "template.pptx")
    monkeypatch.setattr(generate_ppt_module, "PROJECT_ROOT", str(tmp_path))
    return tmp_path


def _write_client_data(project_dir, client_name="Acme"):
    data = {
        "name": client_name,
        "start_date_string": "01/04/2026",
        "end_date_string": "30/04/2026",
        "account_type": "Ecommerce",
        "paid_data_mom": {"Total": {"Cost": {"curr": "£1,000", "prev": "£900", "pct": "+11%"}}},
    }
    with open(project_dir / "storage" / f"{client_name}_monthly_data.json", "w") as f:
        json.dump(data, f)
    return data


class TestGenerateSkeletonPpt:
    def test_raises_if_no_client_data(self, isolated_project):
        with pytest.raises(FileNotFoundError):
            generate_skeleton_ppt("NoSuchClient", {"teams": []})

    def test_persists_checkpoint_and_builds_both_draft_variants(self, isolated_project):
        _write_client_data(isolated_project)
        teams_content = {"teams": [{"team": "ppc", "overviews": [], "plan_json": None}]}
        output_path, presentation_path = generate_skeleton_ppt(
            "Acme", teams_content, output_path=str(isolated_project / "slides" / "skeleton.pptx")
        )
        assert os.path.exists(output_path)
        assert os.path.exists(presentation_path)
        checkpoint_path = isolated_project / "storage" / "Acme_skeleton_content.json"
        assert json.loads(checkpoint_path.read_text()) == teams_content


class TestGenerateFinalPpt:
    def _write_checkpoint(self, project_dir, client_name="Acme"):
        teams_content = {"teams": [{"team": "ppc", "overviews": [], "plan_json": None}]}
        with open(project_dir / "storage" / f"{client_name}_skeleton_content.json", "w") as f:
            json.dump(teams_content, f)

    def test_raises_if_no_skeleton_checkpoint(self, isolated_project):
        _write_client_data(isolated_project)
        with pytest.raises(FileNotFoundError):
            generate_ppt("Acme", slide_content={"trends": []})

    def test_merges_trends_into_ppc_block(self, isolated_project):
        _write_client_data(isolated_project)
        self._write_checkpoint(isolated_project)
        trends = [{"title": "t", "summary": "s", "bullets": [], "graph": {}}]

        output_path, presentation_path, excel_path = generate_ppt("Acme", slide_content={"trends": trends})

        assert os.path.exists(output_path)
        assert os.path.exists(presentation_path)
        saved = json.loads((isolated_project / "storage" / "Acme_monthly_content.json").read_text())
        ppc_block = next(b for b in saved["teams"] if b["team"] == "ppc")
        assert ppc_block["trends"] == trends
